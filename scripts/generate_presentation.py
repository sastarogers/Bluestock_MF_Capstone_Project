"""Generate 12-Slide Presentation
=================================
Build a professional PPTX presentation summarising the
Bluestock MF Capstone Project.

Usage:
    python3 scripts/generate_presentation.py
"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-8s  %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger(__name__)

BASE_DIR = Path(__file__).resolve().parents[1]
REPORTS_DIR = BASE_DIR / "reports"
CHART_DIR_DAY3 = REPORTS_DIR / "charts" / "day3"
CHART_DIR_DAY4 = REPORTS_DIR / "charts" / "day4"
DASHBOARD_DIR = BASE_DIR / "dashboard"

# Colours
NAVY = RGBColor(0x01, 0x29, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xF0, 0x55, 0x37)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
LIGHT_NAVY = RGBColor(0x1E, 0x3A, 0x8A)


def add_bg(slide, color=NAVY):
    """Fill the slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=14, color=WHITE):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf


def safe_add_image(slide, path, left, top, width=None, height=None):
    """Add an image if the file exists."""
    if path.exists():
        kwargs = {"image_file": str(path), "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(**kwargs)
        return True
    return False


def slide_01_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide)

    # Logo
    logo = DASHBOARD_DIR / "bluestock_logo.png"
    safe_add_image(slide, logo, Inches(3.5), Inches(0.5), width=Inches(3))

    add_text_box(slide, Inches(0.5), Inches(3), Inches(9), Inches(1),
                 "Bluestock Mutual Fund Capstone", font_size=32, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(4), Inches(9), Inches(0.6),
                 "Data Engineering & Analytics — Final Presentation",
                 font_size=18, alignment=PP_ALIGN.CENTER)

    # Accent line
    from pptx.shapes.autoshape import Shape
    shape = slide.shapes.add_shape(
        1, Inches(3), Inches(4.8), Inches(4), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(0.5),
                 "Prepared by: Atharva Ranjan  |  June 2026",
                 font_size=14, alignment=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    """Slide 2: Problem & Objective."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "Problem & Objective", font_size=28, bold=True, color=ACCENT)

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(1),
                 "The Indian mutual fund industry manages Rs 81+ lakh crore across 1,900+ schemes. "
                 "Fund selection requires comparing NAV trends, risk metrics, and manager performance "
                 "across complex, fragmented datasets.",
                 font_size=14)

    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(0.6),
                 "Project Objectives", font_size=20, bold=True, color=ACCENT)

    add_bullet_frame(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(3.5), [
        "Build an end-to-end ETL pipeline from raw CSVs to a star-schema warehouse",
        "Clean and validate 10 datasets with 100K+ records",
        "Perform exploratory analysis with 18 publication-quality charts",
        "Compute fund scorecards with CAGR, Sharpe, Alpha, and drawdown metrics",
        "Deliver an interactive 4-page Tableau dashboard",
        "Advanced analytics: VaR/CVaR, rolling Sharpe, investor cohorts, HHI",
    ], font_size=13)


def slide_03_data_sources(prs):
    """Slide 3: Data Sources."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "Data Sources", font_size=28, bold=True, color=ACCENT)

    datasets = [
        "01  Fund Master — 40 schemes, fund houses, categories, risk grades",
        "02  NAV History — 46K daily NAV records → 64K after calendar fill",
        "03  AUM by Fund House — quarterly AUM across 10 fund houses",
        "04  Monthly SIP Inflows — 48 months of SIP trend data",
        "05  Category Inflows — monthly net flows by fund category",
        "06  Industry Folio Count — total folio growth metrics",
        "07  Scheme Performance — returns, alpha, beta, Sharpe for 40 funds",
        "08  Investor Transactions — 50K transactions with demographics",
        "09  Portfolio Holdings — stock-level fund compositions",
        "10  Benchmark Indices — NIFTY50, NIFTY100 daily closes",
    ]

    add_bullet_frame(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(5), datasets,
                     font_size=12)

    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.5),
                 "Live Feed: mfapi.in — 6 bluechip scheme NAVs",
                 font_size=12, bold=True)


def slide_04_architecture(prs):
    """Slide 4: Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "Architecture & ETL Flow", font_size=28, bold=True, color=ACCENT)

    # Architecture boxes
    stages = [
        ("Raw CSVs\n(10 files)", Inches(0.3)),
        ("Python\nCleaning", Inches(2.3)),
        ("Processed\nCSVs", Inches(4.3)),
        ("SQLite\nStar Schema", Inches(6.3)),
        ("Analytics\n& Dashboard", Inches(8.3)),
    ]

    for label, left in stages:
        shape = slide.shapes.add_shape(5, left, Inches(2), Inches(1.5), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.text = label
        for p in tf.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    # Arrow labels
    for x_pos in [Inches(1.8), Inches(3.8), Inches(5.8), Inches(7.8)]:
        add_text_box(slide, x_pos, Inches(2.3), Inches(0.5), Inches(0.5),
                     "→", font_size=24, bold=True, color=ACCENT,
                     alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(4), Inches(9), Inches(0.6),
                 "Star Schema: dim_fund, dim_date → fact_nav, fact_transactions, fact_performance, fact_aum",
                 font_size=13, bold=True)

    add_bullet_frame(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(2.5), [
        "11 SQLite tables with foreign-key integrity",
        "Calendar forward-fill: 46K → 64K NAV rows",
        "Automated quality reports with validation checks",
        "Master pipeline: run_pipeline.py with 10 stages",
    ], font_size=13)


def slide_05_eda_1(prs):
    """Slide 5: EDA Highlights 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "EDA Highlights — NAV & AUM", font_size=28, bold=True, color=ACCENT)

    safe_add_image(slide, CHART_DIR_DAY3 / "01_nav_trend_all_40_schemes.png",
                   Inches(0.3), Inches(1.2), width=Inches(4.5))
    safe_add_image(slide, CHART_DIR_DAY3 / "03_aum_growth_by_fund_house.png",
                   Inches(5), Inches(1.2), width=Inches(4.5))

    add_bullet_frame(slide, Inches(0.3), Inches(5.2), Inches(9), Inches(2), [
        "2023 bull-run and 2024 correction windows clearly visible in NAV trends",
        "SBI Mutual Fund dominates AUM — Rs 12.5 lakh crore reference marker",
    ], font_size=12)


def slide_06_eda_2(prs):
    """Slide 6: EDA Highlights 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "EDA Highlights — SIP & Demographics", font_size=28, bold=True, color=ACCENT)

    safe_add_image(slide, CHART_DIR_DAY3 / "05_sip_inflow_time_series.png",
                   Inches(0.3), Inches(1.2), width=Inches(4.5))
    safe_add_image(slide, CHART_DIR_DAY3 / "09_age_group_distribution_pie.png",
                   Inches(5.2), Inches(1.2), width=Inches(3.5))

    add_bullet_frame(slide, Inches(0.3), Inches(5.2), Inches(9), Inches(2), [
        "SIP inflows reached Rs 31,002 Cr all-time high in December 2025",
        "Investor participation distributed across multiple age bands",
        "Industry folios doubled from 13.26 Cr to 26.12 Cr",
    ], font_size=12)


def slide_07_performance_1(prs):
    """Slide 7: Performance Metrics 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "Performance Metrics — Scorecard", font_size=28, bold=True, color=ACCENT)

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.8),
                 "Composite 0–100 Fund Score = 30% CAGR + 25% Sharpe + 20% Alpha + 15% Expense + 10% Drawdown",
                 font_size=13, bold=True)

    add_bullet_frame(slide, Inches(0.5), Inches(2.2), Inches(9), Inches(4.5), [
        "1-year and 3-year NAV CAGR computed for all 40 funds",
        "Sharpe Ratio using 6.5% annual risk-free rate proxy",
        "Sortino Ratio using downside-only daily volatility",
        "Alpha and Beta via OLS regression against NIFTY100",
        "Maximum drawdown with peak, trough, and recovery dates",
        "Tracking error vs NIFTY50 and NIFTY100 for top 5 funds",
        "5-year CAGR flagged as unavailable (data starts Jan 2022)",
    ], font_size=13)


def slide_08_performance_2(prs):
    """Slide 8: Performance Metrics 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "Benchmark Comparison & Risk", font_size=28, bold=True, color=ACCENT)

    safe_add_image(slide, CHART_DIR_DAY4 / "benchmark_comparison_top5_vs_indices.png",
                   Inches(0.3), Inches(1.2), width=Inches(4.5))

    var_chart = REPORTS_DIR / "var_cvar_chart.png"
    safe_add_image(slide, var_chart, Inches(5), Inches(1.2), width=Inches(4.5))

    add_bullet_frame(slide, Inches(0.3), Inches(5.2), Inches(9), Inches(2), [
        "Top 5 scorecard funds outperformed both NIFTY50 and NIFTY100 benchmarks",
        "VaR (95%) and CVaR computed for all 40 schemes for tail-risk assessment",
    ], font_size=12)


def slide_09_dashboard_1(prs):
    """Slide 9: Dashboard Screenshots 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "Dashboard — Industry & Performance", font_size=28, bold=True, color=ACCENT)

    safe_add_image(slide, DASHBOARD_DIR / "1. Industry Overview.png",
                   Inches(0.3), Inches(1.2), width=Inches(4.5))
    safe_add_image(slide, DASHBOARD_DIR / "2. Fund Performance.png",
                   Inches(5), Inches(1.2), width=Inches(4.5))

    add_text_box(slide, Inches(0.3), Inches(5.5), Inches(4.5), Inches(0.5),
                 "Page 1: Industry Overview", font_size=11, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5), Inches(5.5), Inches(4.5), Inches(0.5),
                 "Page 2: Fund Performance", font_size=11, bold=True,
                 alignment=PP_ALIGN.CENTER)


def slide_10_dashboard_2(prs):
    """Slide 10: Dashboard Screenshots 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "Dashboard — Investors & Market", font_size=28, bold=True, color=ACCENT)

    safe_add_image(slide, DASHBOARD_DIR / "3. Investor Analytics.png",
                   Inches(0.3), Inches(1.2), width=Inches(4.5))
    safe_add_image(slide, DASHBOARD_DIR / "4. SIP & Market Trends.png",
                   Inches(5), Inches(1.2), width=Inches(4.5))

    add_text_box(slide, Inches(0.3), Inches(5.5), Inches(4.5), Inches(0.5),
                 "Page 3: Investor Analytics", font_size=11, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5), Inches(5.5), Inches(4.5), Inches(0.5),
                 "Page 4: SIP & Market Trends", font_size=11, bold=True,
                 alignment=PP_ALIGN.CENTER)


def slide_11_key_findings(prs):
    """Slide 11: Key Findings."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "Key Findings", font_size=28, bold=True, color=ACCENT)

    add_bullet_frame(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5), [
        "SBI Mutual Fund leads with highest AUM; top 3 fund houses control >50% of industry AUM",
        "SIP inflows grew 67% over 4 years, reaching Rs 31,002 Cr in Dec 2025",
        "Industry folios nearly doubled (13.26 Cr → 26.12 Cr), showing strong retail adoption",
        "Equity funds dominate category inflows; debt categories show seasonal volatility",
        "Top-scored funds consistently outperformed NIFTY50 and NIFTY100 benchmarks",
        "Forward-filled NAV data enabled accurate daily return and CAGR calculations",
        "Alpha analysis revealed most equity funds have positive alpha vs NIFTY100",
        "HHI analysis shows moderate sector concentration in most equity portfolios",
    ], font_size=14)


def slide_12_thank_you(prs):
    """Slide 12: Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(2), Inches(9), Inches(1),
                 "Thank You", font_size=40, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Accent line
    shape = slide.shapes.add_shape(
        1, Inches(3), Inches(3.2), Inches(4), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.5),
                 "Bluestock Mutual Fund Capstone Project",
                 font_size=16, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(0.5),
                 "Prepared by: Atharva Ranjan",
                 font_size=14, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(0.5),
                 "GitHub: github.com/sastarogers/Bluestock_MF_Capstone_Project",
                 font_size=12, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(5.8), Inches(9), Inches(0.5),
                 "June 2026", font_size=12, alignment=PP_ALIGN.CENTER)


def main() -> None:
    """Generate the 12-slide presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_data_sources(prs)
    slide_04_architecture(prs)
    slide_05_eda_1(prs)
    slide_06_eda_2(prs)
    slide_07_performance_1(prs)
    slide_08_performance_2(prs)
    slide_09_dashboard_1(prs)
    slide_10_dashboard_2(prs)
    slide_11_key_findings(prs)
    slide_12_thank_you(prs)

    output_path = REPORTS_DIR / "Bluestock_MF_Presentation.pptx"
    prs.save(str(output_path))
    log.info("Presentation generated: %s (%d slides)", output_path, len(prs.slides))


if __name__ == "__main__":
    main()
